import csv
import json
import mimetypes
import os
import sqlite3
import zipfile
from typing import Any, Dict, List, Optional, Tuple

from chatvault_db import add_message

UPLOAD_ROOT = os.path.abspath(os.getenv("CHATVAULT_UPLOAD_DIR", "uploads"))


def _build_text_content(base_content: str, json_patch: Optional[Dict[str, Any]], content: Optional[str], ext: str) -> Tuple[str, str]:
    if json_patch is not None:
        try:
            existing = json.loads(base_content) if base_content.strip() else {}
            if not isinstance(existing, dict):
                raise ValueError("existing content is not a JSON object")
            existing.update(json_patch)
            return json.dumps(existing, ensure_ascii=False, indent=2), "application/json"
        except json.JSONDecodeError:
            raise ValueError("could not parse existing file as JSON")
    content_to_write = content if content is not None else base_content
    content_type = "text/plain" if ext not in {".json"} else "application/json"
    return content_to_write, content_type


def _safe_path(path: str, allow_missing: bool = False) -> str:
    candidate = path
    if not os.path.isabs(candidate):
        candidate = os.path.join(UPLOAD_ROOT, candidate)
    abs_path = os.path.abspath(candidate)
    if not abs_path.startswith(UPLOAD_ROOT):
        raise ValueError("file path is outside upload root")
    if not allow_missing and not os.path.exists(abs_path):
        raise ValueError("file does not exist on disk")
    return abs_path


def _find_by_basename(filename: str) -> Optional[str]:
    target = filename.strip()
    if not target:
        return None
    for root, _, files in os.walk(UPLOAD_ROOT):
        if target in files:
            return os.path.abspath(os.path.join(root, target))
    return None


def _fetch_meta(con: sqlite3.Connection, message_id: int) -> Tuple[Dict[str, Any], Optional[int]]:
    cur = con.cursor()
    cur.execute("SELECT meta_json, conversation_id FROM messages WHERE id=?", (message_id,))
    row = cur.fetchone()
    if not row:
        raise ValueError(f"message_id {message_id} not found")
    raw = row[0] or "{}"
    conv_id = int(row[1]) if row[1] is not None else None
    try:
        return json.loads(raw), conv_id
    except json.JSONDecodeError:
        return {}, conv_id


def resolve_file_path(con: sqlite3.Connection, file_id: str) -> Tuple[str, Dict[str, Any], Optional[int]]:
    fid = (file_id or "").strip()
    if not fid:
        raise ValueError("file_id is required")

    meta: Dict[str, Any] = {}
    conv_id: Optional[int] = None
    if fid.isdigit():
        meta, conv_id = _fetch_meta(con, int(fid))
        stored = meta.get("stored_path")
        if not stored:
            raise ValueError("message has no stored_path")
        return _safe_path(stored), meta, conv_id

    try:
        return _safe_path(fid), meta, conv_id
    except ValueError:
        located = _find_by_basename(fid)
        if located:
            return _safe_path(located), meta, conv_id
        raise


def _log_sidecar(con: sqlite3.Connection, conv_id: Optional[int], path: str, content_type: str) -> Optional[int]:
    if conv_id is None:
        return None
    size_bytes = os.path.getsize(path)
    try:
        return add_message(
            con=con,
            conversation_id=conv_id,
            source="ai_file",
            role="assistant",
            content=f"[AI generated file: {os.path.basename(path)} | {size_bytes} bytes | stored at {path}]",
            meta={
                "filename": os.path.basename(path),
                "bytes": size_bytes,
                "content_type": content_type,
                "stored_path": path,
                "source": "ai_file",
            },
        )
    except Exception:
        return None


def _safe_int(value: Any) -> Optional[int]:
    try:
        return int(value)
    except (TypeError, ValueError):
        return None


def _safe_zip_member(name: str) -> bool:
    if not name:
        return False
    norm = os.path.normpath(name)
    if norm.startswith("..") or os.path.isabs(norm):
        return False
    return True


def _apply_paragraph_replacements(doc: Any, replacements: Optional[List[Dict[str, Any]]]) -> None:
    for item in replacements or []:
        idx = _safe_int(item.get("index"))
        if idx is None:
            continue
        if 0 <= idx < len(doc.paragraphs):
            doc.paragraphs[idx].text = item.get("text", "")


def _apply_paragraph_find_replace(doc: Any, find_replace: Optional[List[Dict[str, Any]]]) -> None:
    for fr in find_replace or []:
        find_text = fr.get("find") if isinstance(fr, dict) else None
        replace_text = fr.get("replace") if isinstance(fr, dict) else None
        if not isinstance(find_text, str) or replace_text is None:
            continue
        for para in doc.paragraphs:
            if para.text and find_text in para.text:
                para.text = para.text.replace(find_text, str(replace_text))


def _append_paragraphs(doc: Any, paragraphs: Optional[List[str]]) -> None:
    for text in paragraphs or []:
        doc.add_paragraph(text)


def modify_text_file(
    con: sqlite3.Connection,
    file_id: str,
    content: Optional[str] = None,
    json_patch: Optional[Dict[str, Any]] = None,
    output_mode: str = "sidecar",
    output_suffix: str = "-edited",
    encoding: str = "utf-8",
) -> Dict[str, Any]:
    if content is None and json_patch is None:
        raise ValueError("content or json_patch is required")

    path, _meta, conv_id = resolve_file_path(con, file_id)
    ext = os.path.splitext(path)[1].lower()

    base_content: str = ""
    if os.path.exists(path):
        with open(path, "r", encoding=encoding, errors="replace") as fh:
            base_content = fh.read()

    out_path = path if output_mode == "replace" else _safe_path(f"{os.path.splitext(path)[0]}{output_suffix}{ext}", allow_missing=True)

    content_to_write, content_type = _build_text_content(base_content, json_patch, content, ext)

    with open(out_path, "w", encoding=encoding) as fh:
        fh.write(content_to_write)

    message_id = _log_sidecar(con, conv_id, out_path, content_type)

    return {
        "ok": True,
        "output_path": out_path,
        "bytes": os.path.getsize(out_path),
        "content_type": content_type,
        "message_id": message_id,
    }


def modify_csv_file(
    con: sqlite3.Connection,
    file_id: str,
    rows: Optional[List[List[Any]]] = None,
    append_rows: Optional[List[List[Any]]] = None,
    output_mode: str = "sidecar",
    output_suffix: str = "-edited",
    newline: str = "",
) -> Dict[str, Any]:
    if rows is None and append_rows is None:
        raise ValueError("rows or append_rows is required")

    path, _meta, conv_id = resolve_file_path(con, file_id)
    ext = os.path.splitext(path)[1].lower()
    if ext not in {".csv", ".tsv"}:
        raise ValueError("only CSV/TSV supported")

    delimiter = "\t" if ext == ".tsv" else ","
    out_path = path if output_mode == "replace" else _safe_path(f"{os.path.splitext(path)[0]}{output_suffix}{ext}", allow_missing=True)

    existing: List[List[Any]] = []
    if os.path.exists(path) and append_rows is not None and rows is None:
        with open(path, "r", newline="") as fh:
            reader = csv.reader(fh, delimiter=delimiter)
            existing = list(reader)

    final_rows: List[List[Any]] = []
    if rows is not None:
        final_rows.extend(rows)
    else:
        final_rows.extend(existing)
    if append_rows:
        final_rows.extend(append_rows)

    with open(out_path, "w", newline=newline) as fh:
        writer = csv.writer(fh, delimiter=delimiter)
        writer.writerows(final_rows)

    message_id = _log_sidecar(con, conv_id, out_path, "text/csv")

    return {
        "ok": True,
        "output_path": out_path,
        "bytes": os.path.getsize(out_path),
        "rows_written": len(final_rows),
        "message_id": message_id,
    }


def modify_docx_file(
    con: sqlite3.Connection,
    file_id: str,
    append_paragraphs: Optional[List[str]] = None,
    replace_paragraphs: Optional[List[Dict[str, Any]]] = None,
    find_replace: Optional[List[Dict[str, Any]]] = None,
    output_mode: str = "sidecar",
    output_suffix: str = "-edited",
) -> Dict[str, Any]:
    try:
        from docx import Document  # type: ignore
    except Exception as exc:  # pragma: no cover
        raise RuntimeError("python-docx not installed") from exc

    if not append_paragraphs and not replace_paragraphs and not find_replace:
        raise ValueError("append_paragraphs or replace_paragraphs or find_replace is required")

    path, _meta, conv_id = resolve_file_path(con, file_id)
    ext = os.path.splitext(path)[1].lower()
    if ext != ".docx":
        raise ValueError("only .docx supported")

    doc = Document(path)

    _apply_paragraph_replacements(doc, replace_paragraphs)
    _apply_paragraph_find_replace(doc, find_replace)
    _append_paragraphs(doc, append_paragraphs)

    out_path = path if output_mode == "replace" else _safe_path(f"{os.path.splitext(path)[0]}{output_suffix}{ext}", allow_missing=True)
    doc.save(out_path)

    message_id = _log_sidecar(con, conv_id, out_path, "application/vnd.openxmlformats-officedocument.wordprocessingml.document")

    return {
        "ok": True,
        "output_path": out_path,
        "bytes": os.path.getsize(out_path),
        "message_id": message_id,
    }


def _update_slide_content(slide: Any, title: Optional[str], body: Optional[str]) -> None:
    if title and slide.shapes.title:
        slide.shapes.title.text = title
    if not body:
        return
    for shape in slide.placeholders:
        if shape != slide.shapes.title and shape.has_text_frame:
            shape.text_frame.text = body
            break


def _replace_slide_content(prs: Any, replace_slides: Optional[List[Dict[str, Any]]]) -> None:
    for item in replace_slides or []:
        idx = _safe_int(item.get("index"))
        if idx is None or not (0 <= idx < len(prs.slides)):
            continue
        _update_slide_content(prs.slides[idx], item.get("title"), item.get("body"))


def _apply_shape_edits(prs: Any, shape_edits: Optional[List[Dict[str, Any]]]) -> None:
    for edit in shape_edits or []:
        slide_idx = _safe_int(edit.get("slide_index"))
        text_value = edit.get("text")
        if slide_idx is None or text_value is None:
            continue
        if not (0 <= slide_idx < len(prs.slides)):
            continue

        slide = prs.slides[slide_idx]
        shapes = slide.shapes
        shape_idx = _safe_int(edit.get("shape_index"))
        candidates = [shapes[shape_idx]] if shape_idx is not None and 0 <= shape_idx < len(shapes) else shapes

        for shape in candidates:
            if getattr(shape, "has_text_frame", False):
                shape.text = str(text_value)
                break


def _add_new_slide(prs: Any, add_slide: bool, title: Optional[str], body: Optional[str]) -> None:
    if not add_slide:
        return
    layout_idx = 1 if len(prs.slide_layouts) > 1 else 0
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
    _update_slide_content(slide, title, body)


def modify_pptx_file(
    con: sqlite3.Connection,
    file_id: str,
    add_slide: bool = True,
    title: Optional[str] = None,
    body: Optional[str] = None,
    replace_slides: Optional[List[Dict[str, Any]]] = None,
    shape_edits: Optional[List[Dict[str, Any]]] = None,
    output_mode: str = "sidecar",
    output_suffix: str = "-edited",
) -> Dict[str, Any]:
    try:
        from pptx import Presentation  # type: ignore
    except Exception as exc:  # pragma: no cover
        raise RuntimeError("python-pptx not installed") from exc

    path, _meta, conv_id = resolve_file_path(con, file_id)
    ext = os.path.splitext(path)[1].lower()
    if ext != ".pptx":
        raise ValueError("only .pptx supported")

    prs = Presentation(path)

    _replace_slide_content(prs, replace_slides)
    _apply_shape_edits(prs, shape_edits)
    _add_new_slide(prs, add_slide, title, body)

    out_path = path if output_mode == "replace" else _safe_path(f"{os.path.splitext(path)[0]}{output_suffix}{ext}", allow_missing=True)
    prs.save(out_path)

    message_id = _log_sidecar(con, conv_id, out_path, "application/vnd.openxmlformats-officedocument.presentationml.presentation")

    return {
        "ok": True,
        "output_path": out_path,
        "bytes": os.path.getsize(out_path),
        "message_id": message_id,
    }


def _create_overlay_page(canvas_module: Any, pdf_reader_cls: Any, page_width: float, page_height: float, text: str, x: float, y: float, font_size: int) -> Any:
    from io import BytesIO

    buf = BytesIO()
    c = canvas_module.Canvas(buf, pagesize=(page_width, page_height))
    c.setFont("Helvetica", font_size)
    for line_idx, line in enumerate(text.split("\n")):
        c.drawString(x, y + (len(text.split("\n")) - line_idx - 1) * (font_size + 2), line)
    c.save()
    buf.seek(0)
    return pdf_reader_cls(buf).pages[0]


def _normalize_overlay(ov: Dict[str, Any]) -> Optional[Dict[str, Any]]:
    idx = _safe_int(ov.get("page_index"))
    if idx is None:
        return None
    text = ov.get("text")
    if not isinstance(text, str):
        return None
    return {
        "page_index": idx,
        "text": text,
        "x": float(ov.get("x", 50)) if ov.get("x") is not None else 50.0,
        "y": float(ov.get("y", 50)) if ov.get("y") is not None else 50.0,
        "font_size": int(ov.get("font_size", 12)) if ov.get("font_size") is not None else 12,
    }


def _group_overlays_by_page(overlays: Optional[List[Dict[str, Any]]]) -> Dict[int, List[Dict[str, Any]]]:
    grouped: Dict[int, List[Dict[str, Any]]] = {}
    for ov in overlays or []:
        normalized = _normalize_overlay(ov)
        if not normalized:
            continue
        grouped.setdefault(normalized["page_index"], []).append(normalized)
    return grouped


def _merge_overlays(writer: Any, reader: Any, overlays: Optional[List[Dict[str, Any]]], canvas_module: Any, pdf_reader_cls: Any) -> None:
    overlay_map = _group_overlays_by_page(overlays)
    for page_index, page in enumerate(reader.pages):
        page_obj = page
        for ov in overlay_map.get(page_index, []):
            overlay_page = _create_overlay_page(
                canvas_module,
                pdf_reader_cls,
                float(page_obj.mediabox.width),
                float(page_obj.mediabox.height),
                ov["text"],
                ov["x"],
                ov["y"],
                ov["font_size"],
            )
            page_obj.merge_page(overlay_page)
        writer.add_page(page_obj)


def _append_text_pages(writer: Any, append_text_pages: Optional[List[str]], canvas_module: Any, pdf_reader_cls: Any) -> None:
    if not append_text_pages:
        return
    for text in append_text_pages:
        from io import BytesIO

        buf = BytesIO()
        c = canvas_module.Canvas(buf)
        c.setFont("Helvetica", 12)
        for line_idx, line in enumerate(str(text).split("\n")):
            c.drawString(50, 750 - line_idx * 16, line)
        c.showPage()
        c.save()
        buf.seek(0)
        new_reader = pdf_reader_cls(buf)
        writer.add_page(new_reader.pages[0])


def _append_pdf_files(writer: Any, con: sqlite3.Connection, append_pdf_ids: Optional[List[str]], pdf_reader_cls: Any) -> None:
    if not append_pdf_ids:
        return
    for other_id in append_pdf_ids:
        other_path, _other_meta, _other_conv = resolve_file_path(con, other_id)
        if not other_path.lower().endswith(".pdf"):
            continue
        other_reader = pdf_reader_cls(other_path)
        for p in other_reader.pages:
            writer.add_page(p)


def modify_pdf_file(
    con: sqlite3.Connection,
    file_id: str,
    append_pdf_ids: Optional[List[str]] = None,
    append_text_pages: Optional[List[str]] = None,
    overlay_texts: Optional[List[Dict[str, Any]]] = None,
    output_mode: str = "sidecar",
    output_suffix: str = "-edited",
) -> Dict[str, Any]:
    try:
        from PyPDF2 import PdfReader, PdfWriter  # type: ignore
    except Exception as exc:  # pragma: no cover
        raise RuntimeError("PyPDF2 not installed") from exc

    try:
        from reportlab.pdfgen import canvas  # type: ignore
    except Exception as exc:  # pragma: no cover
        raise RuntimeError("reportlab not installed") from exc

    base_path, _meta, conv_id = resolve_file_path(con, file_id)
    if not base_path.lower().endswith(".pdf"):
        raise ValueError("only .pdf supported")

    reader = PdfReader(base_path)
    writer = PdfWriter()

    _merge_overlays(writer, reader, overlay_texts, canvas, PdfReader)
    _append_text_pages(writer, append_text_pages, canvas, PdfReader)
    _append_pdf_files(writer, con, append_pdf_ids, PdfReader)

    ext = ".pdf"
    out_path = base_path if output_mode == "replace" else _safe_path(f"{os.path.splitext(base_path)[0]}{output_suffix}{ext}", allow_missing=True)
    with open(out_path, "wb") as fh:
        writer.write(fh)

    message_id = _log_sidecar(con, conv_id, out_path, "application/pdf")

    return {
        "ok": True,
        "output_path": out_path,
        "bytes": os.path.getsize(out_path),
        "message_id": message_id,
    }


def inspect_zip_file(
    con: sqlite3.Connection,
    file_id: str,
    max_entries: Optional[int] = None,
) -> Dict[str, Any]:
    path, _meta, _conv_id = resolve_file_path(con, file_id)
    if not path.lower().endswith(".zip"):
        raise ValueError("only .zip supported")
    if not zipfile.is_zipfile(path):
        raise ValueError("not a valid zip archive")

    entries: List[Dict[str, Any]] = []
    with zipfile.ZipFile(path, "r") as zf:
        infos = zf.infolist()
        total_entries = len(infos)
        cap = _safe_int(max_entries) if max_entries is not None else None
        limit = cap if cap and cap > 0 else total_entries
        for info in infos[:limit]:
            entries.append(
                {
                    "name": info.filename,
                    "size": info.file_size,
                    "compressed_size": info.compress_size,
                    "is_dir": info.is_dir(),
                }
            )
    return {
        "ok": True,
        "entries": entries,
        "count": total_entries,
        "truncated": cap is not None and total_entries > len(entries),
    }


def extract_zip_file(
    con: sqlite3.Connection,
    file_id: str,
    members: Optional[List[str]] = None,
    output_suffix: str = "-unzipped",
) -> Dict[str, Any]:
    path, _meta, conv_id = resolve_file_path(con, file_id)
    if not path.lower().endswith(".zip"):
        raise ValueError("only .zip supported")
    if not zipfile.is_zipfile(path):
        raise ValueError("not a valid zip archive")

    base_dir = os.path.dirname(path)
    zip_stem = os.path.splitext(os.path.basename(path))[0]
    dest_base = _safe_path(os.path.join(base_dir, f"{zip_stem}{output_suffix}"), allow_missing=True)
    dest_dir = dest_base
    counter = 1
    while os.path.exists(dest_dir):
        dest_dir = _safe_path(f"{dest_base}-{counter}", allow_missing=True)
        counter += 1

    os.makedirs(dest_dir, exist_ok=True)

    extracted: List[Dict[str, Any]] = []
    skipped: List[Dict[str, Any]] = []
    with zipfile.ZipFile(path, "r") as zf:
        target_members = members or [info.filename for info in zf.infolist()]
        for name in target_members:
            if not _safe_zip_member(name):
                skipped.append({"name": name, "reason": "unsafe_path"})
                continue
            try:
                info = zf.getinfo(name)
            except KeyError:
                skipped.append({"name": name, "reason": "not_found"})
                continue

            if info.is_dir():
                dir_path = _safe_path(os.path.join(dest_dir, os.path.normpath(info.filename)), allow_missing=True)
                os.makedirs(dir_path, exist_ok=True)
                continue

            data = zf.read(info)
            target_path = _safe_path(os.path.join(dest_dir, os.path.normpath(info.filename)), allow_missing=True)
            os.makedirs(os.path.dirname(target_path), exist_ok=True)
            with open(target_path, "wb") as fh:
                fh.write(data)

            content_type = mimetypes.guess_type(info.filename)[0] or "application/octet-stream"
            message_id = _log_sidecar(con, conv_id, target_path, content_type)
            extracted.append(
                {
                    "name": info.filename,
                    "output_path": target_path,
                    "bytes": len(data),
                    "message_id": message_id,
                }
            )

    return {
        "ok": True,
        "output_dir": dest_dir,
        "extracted": extracted,
        "skipped": skipped,
    }


def create_zip_archive(
    con: sqlite3.Connection,
    file_ids: List[str],
    zip_filename: Optional[str] = None,
) -> Dict[str, Any]:
    if not file_ids:
        raise ValueError("file_ids is required")

    resolved: List[str] = []
    conv_id: Optional[int] = None
    for fid in file_ids:
        p, _meta, maybe_conv = resolve_file_path(con, fid)
        resolved.append(p)
        if conv_id is None and maybe_conv is not None:
            conv_id = maybe_conv

    base_dir = os.path.dirname(resolved[0])
    name = (zip_filename or "ai-archive.zip").strip() or "ai-archive.zip"
    if not name.lower().endswith(".zip"):
        name = f"{name}.zip"

    target = _safe_path(os.path.join(base_dir, name), allow_missing=True)
    base_target = target
    counter = 1
    while os.path.exists(target):
        target = _safe_path(f"{os.path.splitext(base_target)[0]}-{counter}.zip", allow_missing=True)
        counter += 1

    common_prefix = os.path.commonpath(resolved)
    with zipfile.ZipFile(target, "w", compression=zipfile.ZIP_DEFLATED) as zf:
        for path in resolved:
            arcname = os.path.relpath(path, common_prefix)
            zf.write(path, arcname=arcname)

    message_id = _log_sidecar(con, conv_id, target, "application/zip")

    return {
        "ok": True,
        "zip_path": target,
        "files_added": len(resolved),
        "message_id": message_id,
    }
