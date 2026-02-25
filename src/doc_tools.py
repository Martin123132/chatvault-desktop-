import os
import json
import sqlite3
from typing import Any, List, Dict, TypeVar

from excel_tools import resolve_file_path

try:  # Optional imports are validated at call time
    from PyPDF2 import PdfReader  # type: ignore
except Exception:  # pragma: no cover - defer missing dep to runtime error
    PdfReader = None  # type: ignore

try:
    from docx import Document  # type: ignore
except Exception:  # pragma: no cover
    Document = None  # type: ignore

try:
    from pptx import Presentation  # type: ignore
except Exception:  # pragma: no cover
    Presentation = None  # type: ignore

DEFAULT_MAX_PAGES = int(os.getenv("CHATVAULT_PDF_MAX_PAGES", "20"))
DEFAULT_MAX_PARAGRAPHS = int(os.getenv("CHATVAULT_DOCX_MAX_PARAGRAPHS", "400"))
DEFAULT_MAX_SLIDES = int(os.getenv("CHATVAULT_PPTX_MAX_SLIDES", "60"))
DEFAULT_MAX_CHARS = int(os.getenv("CHATVAULT_DOC_MAX_CHARS", "12000"))


def _truncate(text: str, cap: int) -> tuple[str, bool]:
    if cap <= 0:
        return text, False
    if len(text) <= cap:
        return text, False
    return text[:cap] + "... [truncated]", True


T = TypeVar("T")


def _assert_available(lib: T | None, name: str) -> T:
    if lib is None:
        raise RuntimeError(f"{name} support not installed")
    return lib


def inspect_pdf(
    con: sqlite3.Connection,
    file_id: str,
    max_pages: int | None = None,
    max_chars: int | None = None,
) -> Dict[str, Any]:
    reader_cls = _assert_available(PdfReader, "PDF")
    path, meta, _ = resolve_file_path(con, file_id)
    if not path.lower().endswith(".pdf"):
        raise ValueError("only PDF files are supported")

    page_cap = max_pages or DEFAULT_MAX_PAGES
    char_cap = max_chars or DEFAULT_MAX_CHARS

    with open(path, "rb") as fh:
        reader = reader_cls(fh)
        pages: List[Dict[str, Any]] = []
        for idx, page in enumerate(reader.pages):
            if idx >= page_cap:
                break
            try:
                raw_text = page.extract_text() or ""
            except Exception:
                raw_text = ""
            raw_text = raw_text.strip()
            text, truncated = _truncate(raw_text, char_cap)
            pages.append({
                "page_index": idx,
                "text": text,
                "chars": len(raw_text),
                "truncated": truncated,
            })

    return {
        "ok": True,
        "file_path": path,
        "page_count": len(getattr(reader, "pages", [])),
        "pages": pages,
        "meta": meta,
    }


def inspect_docx(
    con: sqlite3.Connection,
    file_id: str,
    max_paragraphs: int | None = None,
    max_chars: int | None = None,
) -> Dict[str, Any]:
    doc_cls = _assert_available(Document, "DOCX")
    path, meta, _ = resolve_file_path(con, file_id)
    if not path.lower().endswith(".docx"):
        raise ValueError("only .docx files are supported")

    para_cap = max_paragraphs or DEFAULT_MAX_PARAGRAPHS
    char_cap = max_chars or DEFAULT_MAX_CHARS

    doc = doc_cls(path)
    paragraphs = []
    for idx, p in enumerate(doc.paragraphs):
        if idx >= para_cap:
            break
        raw_text = (p.text or "").strip()
        text, truncated = _truncate(raw_text, char_cap)
        paragraphs.append({
            "index": idx,
            "text": text,
            "chars": len(raw_text),
            "style": getattr(p.style, "name", None),
            "truncated": truncated,
        })

    return {
        "ok": True,
        "file_path": path,
        "paragraph_count": len(doc.paragraphs),
        "paragraphs": paragraphs,
        "meta": meta,
    }


def _collect_slide_texts(slide: Any) -> List[str]:
    texts: List[str] = []
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False) and shape.text_frame:
            txt = shape.text_frame.text or ""
            if txt.strip():
                texts.append(txt.strip())
    return texts


def inspect_pptx(
    con: sqlite3.Connection,
    file_id: str,
    max_slides: int | None = None,
    max_chars: int | None = None,
) -> Dict[str, Any]:
    pres_cls = _assert_available(Presentation, "PPTX")
    path, meta, _ = resolve_file_path(con, file_id)
    if not path.lower().endswith(".pptx"):
        raise ValueError("only .pptx files are supported")

    slide_cap = max_slides or DEFAULT_MAX_SLIDES
    char_cap = max_chars or DEFAULT_MAX_CHARS

    pres = pres_cls(path)
    slides_payload = []
    for idx, slide in enumerate(pres.slides):
        if idx >= slide_cap:
            break
        texts = _collect_slide_texts(slide)
        combined = "\n\n".join(texts).strip()
        text, truncated = _truncate(combined, char_cap)
        slides_payload.append({
            "index": idx,
            "text": text,
            "chars": len(combined),
            "truncated": truncated,
        })

    return {
        "ok": True,
        "file_path": path,
        "slide_count": len(pres.slides),
        "slides": slides_payload,
        "meta": meta,
    }
